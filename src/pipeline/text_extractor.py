"""Text extraction from various file formats."""

from pathlib import Path
from typing import Optional

from src.config import Config
from src.models.file_descriptor import FileDescriptor
from src.utils.helpers import safe_read_text, is_text_extension


class TextExtractor:
    def __init__(self, config: Optional[Config] = None):
        self.config = config

    def extract(self, file_node: FileDescriptor) -> str:
        path = Path(file_node.path)
        ext = file_node.extension.lower()
        if not path.is_file():
            file_node.has_text = False
            return ""

        text = ""
        try:
            if ext == ".pdf":
                text = self._extract_pdf(path)
            elif ext in (".docx",):
                text = self._extract_docx(path)
            elif ext in (".pptx",):
                text = self._extract_pptx(path)
            elif is_text_extension(ext) or ext in (".doc",):
                text = safe_read_text(path)
            else:
                text = safe_read_text(path, max_bytes=64_000)
        except Exception:
            text = ""

        file_node.content_text = text
        file_node.has_text = bool(text.strip())
        return text

    def _extract_pdf(self, path: Path) -> str:
        import fitz
        doc = fitz.open(str(path))
        parts = []
        for page in doc:
            parts.append(page.get_text())
        doc.close()
        return "\n".join(parts)

    def _extract_docx(self, path: Path) -> str:
        from docx import Document
        doc = Document(str(path))
        return "\n".join(p.text for p in doc.paragraphs)

    def _extract_pptx(self, path: Path) -> str:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    parts.append(shape.text)
        return "\n".join(parts)
